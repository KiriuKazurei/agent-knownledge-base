import hashlib
import html
import json
import mimetypes
import re
import uuid
from pathlib import Path


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".go", ".py", ".js", ".jsx",
    ".ts", ".tsx", ".java", ".cs", ".rs", ".cpp", ".c", ".h", ".json",
    ".yaml", ".yml", ".toml", ".xml", ".sql", ".sh", ".ps1", ".css",
}


def _read_text(path):
    raw = path.read_bytes()
    try:
        return raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        pass
    for encoding in ("gb18030", "utf-16"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            pass
    return raw.decode("utf-8", errors="replace")


def _text_blocks(path):
    text = _read_text(path)
    blocks = []
    current_heading = None
    buffer = []
    start_line = 1
    for line_number, line in enumerate(text.splitlines(), 1):
        heading = re.match(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$", line)
        if heading and buffer:
            blocks.append(("\n".join(buffer), {"kind": "text", "heading": current_heading, "lineStart": start_line, "lineEnd": line_number - 1}))
            buffer = []
        if heading:
            current_heading = heading.group(2)
            start_line = line_number
        if not buffer:
            start_line = line_number
        buffer.append(line)
    if buffer:
        blocks.append(("\n".join(buffer), {"kind": "text", "heading": current_heading, "lineStart": start_line, "lineEnd": len(text.splitlines()) or 1}))
    return blocks


def _pdf_blocks(path):
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return [((page.extract_text() or "").strip(), {"kind": "pdf", "page": index + 1}) for index, page in enumerate(reader.pages)]


def _docx_blocks(path):
    from docx import Document
    document = Document(str(path))
    blocks = []
    heading = None
    for index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name if paragraph.style else ""
        if style.lower().startswith("heading"):
            heading = text
        blocks.append((text, {"kind": "docx", "paragraph": index + 1, "heading": heading, "style": style}))
    for table_index, table in enumerate(document.tables):
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        blocks.append(("\n".join(rows), {"kind": "docx-table", "table": table_index + 1}))
    return blocks


def _xlsx_blocks(path):
    from openpyxl import load_workbook
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    blocks = []
    for sheet in workbook.worksheets:
        rows = []
        start = 1
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
            values = ["" if value is None else str(value) for value in row]
            if not any(values):
                continue
            rows.append("\t".join(values))
            if len(rows) >= 60:
                blocks.append(("\n".join(rows), {"kind": "xlsx", "sheet": sheet.title, "rowStart": start, "rowEnd": row_index}))
                rows, start = [], row_index + 1
        if rows:
            blocks.append(("\n".join(rows), {"kind": "xlsx", "sheet": sheet.title, "rowStart": start, "rowEnd": sheet.max_row}))
    workbook.close()
    return blocks


def _pptx_blocks(path):
    from pptx import Presentation
    presentation = Presentation(str(path))
    blocks = []
    for slide_index, slide in enumerate(presentation.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        blocks.append(("\n".join(texts), {"kind": "pptx", "slide": slide_index + 1}))
    return blocks


def _html_blocks(path):
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(_read_text(path), "html.parser")
        for node in soup(["script", "style", "noscript", "template"]):
            node.decompose()
        title = soup.title.get_text(" ", strip=True) if soup.title else path.stem
        blocks = []
        heading = None
        for node in soup.find_all(["h1", "h2", "h3", "p", "li", "pre", "table"]):
            text = node.get_text(" ", strip=True)
            if not text:
                continue
            if node.name in {"h1", "h2", "h3"}:
                heading = text
            blocks.append((text, {"kind": "html", "heading": heading, "element": node.name, "title": title}))
        return blocks
    except ImportError:
        text = re.sub(r"<script[\s\S]*?</script>", " ", _read_text(path), flags=re.I)
        text = html.unescape(re.sub(r"<[^>]+>", " ", text))
        return [(re.sub(r"\s+", " ", text), {"kind": "html"})]


def _split(text, location, target=1200, overlap=120):
    text = text.strip()
    if not text:
        return []
    if len(text) <= target:
        return [(text, dict(location))]
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + target)
        if end < len(text):
            boundary = max(text.rfind("\n", start + target // 2, end), text.rfind("。", start + target // 2, end), text.rfind(". ", start + target // 2, end))
            if boundary > start:
                end = boundary + 1
        part = text[start:end].strip()
        if part:
            loc = dict(location)
            loc["charStart"], loc["charEnd"] = start, end
            chunks.append((part, loc))
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return chunks


def parse_document(path_value, document_id, media_type="", original_name=""):
    path = Path(path_value).resolve(strict=True)
    suffix = (Path(original_name).suffix or path.suffix).lower()
    if suffix in TEXT_EXTENSIONS or media_type.startswith("text/"):
        blocks = _text_blocks(path)
    elif suffix == ".pdf":
        blocks = _pdf_blocks(path)
    elif suffix == ".docx":
        blocks = _docx_blocks(path)
    elif suffix == ".xlsx":
        blocks = _xlsx_blocks(path)
    elif suffix == ".pptx":
        blocks = _pptx_blocks(path)
    elif suffix in {".html", ".htm"} or "html" in media_type:
        blocks = _html_blocks(path)
    else:
        raise ValueError("Unsupported document format: {}".format(suffix or media_type))

    chunks = []
    ordinal = 0
    for text, location in blocks:
        for part, part_location in _split(text, location):
            digest = hashlib.sha256(part.encode("utf-8")).hexdigest()
            stable = uuid.uuid5(uuid.NAMESPACE_URL, "{}:{}:{}".format(document_id, ordinal, digest))
            chunks.append({"id": str(stable), "text": part, "location": part_location, "contentHash": digest})
            ordinal += 1
    media = media_type or mimetypes.guess_type(original_name or path.name)[0] or "application/octet-stream"
    return {"title": path.stem, "mediaType": media, "chunks": chunks}
