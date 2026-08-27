import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).parents[1]
sys.path.insert(0, str(ROOT / 'src'))

from knowledge_worker.parsers import parse_document


class OfficePdfTests(unittest.TestCase):
    def test_chinese_docx_extracts_paragraph_and_table(self):
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx is not installed")
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "中文文档.docx"
            document = Document()
            document.add_heading("中文文档标题", level=1)
            document.add_paragraph("中文正文必须保持可检索且不能出现乱码。")
            table = document.add_table(rows=1, cols=2)
            table.rows[0].cells[0].text = "字段"
            table.rows[0].cells[1].text = "中文值"
            document.save(str(source))
            result = parse_document(str(source), "document-docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            text = "\n".join(chunk["text"] for chunk in result["chunks"])
            self.assertIn("中文正文必须保持可检索且不能出现乱码。", text)
            self.assertIn("字段\t中文值", text)
            self.assertTrue(any(chunk["location"].get("kind") == "docx" for chunk in result["chunks"]))
            self.assertNotIn("\ufffd", text)

    def test_chinese_xlsx_extracts_sheet_and_cells(self):
        try:
            from openpyxl import Workbook
        except ImportError:
            self.skipTest("openpyxl is not installed")
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "中文表格.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "中文数据"
            sheet.append(["标题", "内容"])
            sheet.append(["中文记录", "表格内容不能乱码"])
            workbook.save(str(source))
            result = parse_document(str(source), "document-xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            text = "\n".join(chunk["text"] for chunk in result["chunks"])
            self.assertIn("中文记录\t表格内容不能乱码", text)
            self.assertEqual(result["chunks"][0]["location"]["sheet"], "中文数据")
            self.assertNotIn("\ufffd", text)

    def test_chinese_pptx_extracts_slide_text(self):
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx is not installed")
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "中文演示.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "中文演示标题"
            slide.placeholders[1].text = "演示正文必须保留中文。"
            presentation.save(str(source))
            result = parse_document(str(source), "document-pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            text = "\n".join(chunk["text"] for chunk in result["chunks"])
            self.assertIn("中文演示标题", text)
            self.assertIn("演示正文必须保留中文。", text)
            self.assertEqual(result["chunks"][0]["location"]["slide"], 1)
            self.assertNotIn("\ufffd", text)

    def test_pdf_with_unmapped_cjk_fails_closed(self):
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.pdfgen.canvas import Canvas
        except ImportError:
            self.skipTest("reportlab is not installed for PDF fixture generation")
        fonts = [
            Path("C:/Windows/Fonts/simsunb.ttf"),
            Path("C:/Windows/Fonts/Deng.ttf"),
            Path("C:/Windows/Fonts/NotoSansSC-VF.ttf"),
        ]
        font_path = next((path for path in fonts if path.exists()), None)
        if font_path is None:
            self.skipTest("no CJK TrueType font is available")
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "中文报告.pdf"
            pdfmetrics.registerFont(TTFont("KAH-CJK", str(font_path)))
            canvas = Canvas(str(source))
            canvas.setFont("KAH-CJK", 16)
            canvas.drawString(72, 720, "中文 PDF 页面内容不能乱码。")
            canvas.save()
            with self.assertRaises(UnicodeError):
                parse_document(str(source), "document-pdf", "application/pdf")
if __name__ == "__main__":
    unittest.main()
